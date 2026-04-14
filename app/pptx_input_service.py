"""アップロードされた PPTX からテキストを抽出（資料をインプットとして AI に渡す）。"""

from io import BytesIO

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract_text_from_pptx(content: bytes) -> str:
    """スライド内のテキストボックス・表セルからプレーンテキストを抽出。"""
    prs = Presentation(BytesIO(content))
    parts: list[str] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        block_lines: list[str] = []

        for shape in slide.shapes:
            block_lines.extend(_lines_from_shape(shape))

        if block_lines:
            parts.append(
                f"--- スライド {slide_idx} ---\n" + "\n".join(block_lines)
            )

    if not parts:
        return "（この PPTX からはテキストを抽出できませんでした。図形のみのスライドの可能性があります。）"

    return "\n\n".join(parts)


def _lines_from_shape(shape) -> list[str]:
    out: list[str] = []
    try:
        if getattr(shape, "has_text_frame", False) and shape.text_frame:
            for para in shape.text_frame.paragraphs:
                t = (para.text or "").strip()
                if t:
                    out.append(t)

        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.TABLE:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        for para in cell.text_frame.paragraphs:
                            t = (para.text or "").strip()
                            if t:
                                out.append(t)

        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            for sub in shape.shapes:
                out.extend(_lines_from_shape(sub))
    except Exception:
        pass
    return out

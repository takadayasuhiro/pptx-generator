import os
import textwrap
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
from lxml import etree
from app.models import PresentationData, SlideData, ChartData
from app.image_service import fetch_image
from app.config import OUTPUT_DIR

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

THEMES = {
    "ocean_blue": {
        "primary": RGBColor(0x0D, 0x47, 0xA1),
        "accent": RGBColor(0x29, 0xB6, 0xF6),
        "accent2": RGBColor(0x01, 0x57, 0x9B),
        "bg_dark": RGBColor(0x0D, 0x47, 0xA1),
        "bg_light": RGBColor(0xF5, 0xF9, 0xFF),
        "text_dark": RGBColor(0x26, 0x32, 0x38),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "bar": RGBColor(0x29, 0xB6, 0xF6),
        "chart_colors": [
            RGBColor(0x29, 0xB6, 0xF6), RGBColor(0x0D, 0x47, 0xA1),
            RGBColor(0x01, 0x57, 0x9B), RGBColor(0x4F, 0xC3, 0xF7),
            RGBColor(0x03, 0x9B, 0xE5), RGBColor(0x81, 0xD4, 0xFA),
            RGBColor(0x00, 0x83, 0x8F), RGBColor(0x00, 0x60, 0x64),
        ],
        "font_title": "Yu Gothic UI",
        "font_body": "Yu Gothic UI",
        "is_dark": False,
    },
    "midnight": {
        "primary": RGBColor(0x1A, 0x1A, 0x2E),
        "accent": RGBColor(0xE9, 0x4D, 0x6B),
        "accent2": RGBColor(0x53, 0x35, 0x83),
        "bg_dark": RGBColor(0x16, 0x21, 0x3E),
        "bg_light": RGBColor(0x1A, 0x1A, 0x2E),
        "text_dark": RGBColor(0xEE, 0xEE, 0xEE),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "bar": RGBColor(0xE9, 0x4D, 0x6B),
        "chart_colors": [
            RGBColor(0xE9, 0x4D, 0x6B), RGBColor(0x29, 0xB6, 0xF6),
            RGBColor(0xFF, 0xAB, 0x40), RGBColor(0x66, 0xBB, 0x6A),
            RGBColor(0xCE, 0x93, 0xD8), RGBColor(0x4D, 0xD0, 0xE1),
            RGBColor(0xFF, 0xD5, 0x4F), RGBColor(0xA5, 0xD6, 0xA7),
        ],
        "font_title": "Yu Gothic UI",
        "font_body": "Yu Gothic UI",
        "is_dark": True,
    },
    "forest": {
        "primary": RGBColor(0x1B, 0x5E, 0x20),
        "accent": RGBColor(0x66, 0xBB, 0x6A),
        "accent2": RGBColor(0x2E, 0x7D, 0x32),
        "bg_dark": RGBColor(0x1B, 0x5E, 0x20),
        "bg_light": RGBColor(0xF1, 0xF8, 0xE9),
        "text_dark": RGBColor(0x26, 0x32, 0x38),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "bar": RGBColor(0x66, 0xBB, 0x6A),
        "chart_colors": [
            RGBColor(0x2E, 0x7D, 0x32), RGBColor(0x43, 0xA0, 0x47),
            RGBColor(0x66, 0xBB, 0x6A), RGBColor(0x1B, 0x5E, 0x20),
            RGBColor(0x81, 0xC7, 0x84), RGBColor(0x4C, 0xAF, 0x50),
            RGBColor(0x00, 0x96, 0x88), RGBColor(0x00, 0x79, 0x6B),
        ],
        "font_title": "Yu Gothic UI",
        "font_body": "Yu Gothic UI",
        "is_dark": False,
    },
    "sunset": {
        "primary": RGBColor(0xBF, 0x36, 0x0C),
        "accent": RGBColor(0xFF, 0x8F, 0x00),
        "accent2": RGBColor(0xE6, 0x51, 0x00),
        "bg_dark": RGBColor(0xBF, 0x36, 0x0C),
        "bg_light": RGBColor(0xFF, 0xF8, 0xF0),
        "text_dark": RGBColor(0x33, 0x25, 0x1F),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "bar": RGBColor(0xFF, 0x8F, 0x00),
        "chart_colors": [
            RGBColor(0xFF, 0x8F, 0x00), RGBColor(0xBF, 0x36, 0x0C),
            RGBColor(0xE6, 0x51, 0x00), RGBColor(0xFF, 0xB7, 0x4D),
            RGBColor(0xFF, 0x6D, 0x00), RGBColor(0xD8, 0x43, 0x15),
            RGBColor(0xFF, 0xA0, 0x00), RGBColor(0xE5, 0x73, 0x73),
        ],
        "font_title": "Yu Gothic UI",
        "font_body": "Yu Gothic UI",
        "is_dark": False,
    },
    "lavender": {
        "primary": RGBColor(0x4A, 0x14, 0x8C),
        "accent": RGBColor(0xCE, 0x93, 0xD8),
        "accent2": RGBColor(0x7B, 0x1F, 0xA2),
        "bg_dark": RGBColor(0x4A, 0x14, 0x8C),
        "bg_light": RGBColor(0xF3, 0xE5, 0xF5),
        "text_dark": RGBColor(0x2C, 0x2C, 0x2C),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "bar": RGBColor(0xCE, 0x93, 0xD8),
        "chart_colors": [
            RGBColor(0x7B, 0x1F, 0xA2), RGBColor(0xCE, 0x93, 0xD8),
            RGBColor(0x4A, 0x14, 0x8C), RGBColor(0xE1, 0xBE, 0xE7),
            RGBColor(0x9C, 0x27, 0xB0), RGBColor(0xBA, 0x68, 0xC8),
            RGBColor(0x6A, 0x1B, 0x9A), RGBColor(0xAB, 0x47, 0xBC),
        ],
        "font_title": "Yu Gothic UI",
        "font_body": "Yu Gothic UI",
        "is_dark": False,
    },
    "monochrome": {
        "primary": RGBColor(0x21, 0x21, 0x21),
        "accent": RGBColor(0x75, 0x75, 0x75),
        "accent2": RGBColor(0x42, 0x42, 0x42),
        "bg_dark": RGBColor(0x21, 0x21, 0x21),
        "bg_light": RGBColor(0xFA, 0xFA, 0xFA),
        "text_dark": RGBColor(0x21, 0x21, 0x21),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "bar": RGBColor(0x61, 0x61, 0x61),
        "chart_colors": [
            RGBColor(0x26, 0x32, 0x38), RGBColor(0x54, 0x6E, 0x7A),
            RGBColor(0x78, 0x90, 0x9C), RGBColor(0x37, 0x47, 0x4F),
            RGBColor(0x90, 0xA4, 0xAE), RGBColor(0x60, 0x7D, 0x8B),
            RGBColor(0x45, 0x5A, 0x64), RGBColor(0xB0, 0xBE, 0xC5),
        ],
        "font_title": "Yu Gothic UI",
        "font_body": "Yu Gothic UI",
        "is_dark": False,
    },
    "coral": {
        "primary": RGBColor(0xC6, 0x28, 0x28),
        "accent": RGBColor(0xFF, 0x70, 0x43),
        "accent2": RGBColor(0xD8, 0x43, 0x15),
        "bg_dark": RGBColor(0xC6, 0x28, 0x28),
        "bg_light": RGBColor(0xFF, 0xF5, 0xF2),
        "text_dark": RGBColor(0x2C, 0x2C, 0x2C),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "bar": RGBColor(0xFF, 0x70, 0x43),
        "chart_colors": [
            RGBColor(0xEF, 0x53, 0x50), RGBColor(0xFF, 0x70, 0x43),
            RGBColor(0xC6, 0x28, 0x28), RGBColor(0xFF, 0xA2, 0x6D),
            RGBColor(0xD8, 0x43, 0x15), RGBColor(0xFF, 0x8A, 0x65),
            RGBColor(0xE5, 0x39, 0x35), RGBColor(0xFF, 0xAB, 0x91),
        ],
        "font_title": "Yu Gothic UI",
        "font_body": "Yu Gothic UI",
        "is_dark": False,
    },
    "emerald": {
        "primary": RGBColor(0x00, 0x4D, 0x40),
        "accent": RGBColor(0x26, 0xA6, 0x9A),
        "accent2": RGBColor(0x00, 0x79, 0x6B),
        "bg_dark": RGBColor(0x00, 0x4D, 0x40),
        "bg_light": RGBColor(0xE0, 0xF2, 0xF1),
        "text_dark": RGBColor(0x26, 0x32, 0x38),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "bar": RGBColor(0x26, 0xA6, 0x9A),
        "chart_colors": [
            RGBColor(0x00, 0x79, 0x6B), RGBColor(0x26, 0xA6, 0x9A),
            RGBColor(0x00, 0x4D, 0x40), RGBColor(0x80, 0xCB, 0xC4),
            RGBColor(0x00, 0x96, 0x88), RGBColor(0x4D, 0xB6, 0xAC),
            RGBColor(0x00, 0xBF, 0xA5), RGBColor(0xA7, 0xFF, 0xEB),
        ],
        "font_title": "Yu Gothic UI",
        "font_body": "Yu Gothic UI",
        "is_dark": False,
    },
}


def _t(theme_name: str) -> dict:
    return THEMES.get(theme_name, THEMES["ocean_blue"])


def _is_dark(theme_name: str) -> bool:
    return _t(theme_name).get("is_dark", False)


def _bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, left, top, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        sp_pr = shape._element.spPr
        solid_fill = sp_pr.find(qn("a:solidFill"))
        if solid_fill is not None:
            srgb = solid_fill.find(qn("a:srgbClr"))
            if srgb is not None:
                a_alpha = etree.SubElement(srgb, qn("a:alpha"))
                a_alpha.set("val", str(int(alpha * 1000)))
    return shape


def _add_shadow(shape, blur=Pt(6), offset=Pt(3)):
    sp = shape._element
    e_sp_pr = sp.find(qn("p:spPr"))
    if e_sp_pr is None:
        return
    eff = e_sp_pr.find(qn("a:effectLst"))
    if eff is None:
        eff = etree.SubElement(e_sp_pr, qn("a:effectLst"))
    outer = etree.SubElement(eff, qn("a:outerShdw"))
    outer.set("blurRad", str(int(blur)))
    outer.set("dist", str(int(offset)))
    outer.set("dir", "5400000")
    outer.set("algn", "tl")
    outer.set("rotWithShape", "0")
    srgb = etree.SubElement(outer, qn("a:srgbClr"))
    srgb.set("val", "000000")
    a_alpha = etree.SubElement(srgb, qn("a:alpha"))
    a_alpha.set("val", "28000")


def _text(tf, text, font, size, color, bold=False, align=PP_ALIGN.LEFT):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align


def _add_multiline(tf, lines, font, size, color, bold=False,
                   align=PP_ALIGN.LEFT, spacing_before=Pt(4), spacing_after=Pt(4)):
    """Add multiple lines/paragraphs to a text frame."""
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = font
        p.font.size = size
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        p.space_before = spacing_before
        p.space_after = spacing_after


def _add_page_number(slide, num, theme):
    t = _t(theme)
    pg_color = RGBColor(0x77, 0x77, 0x77) if not _is_dark(theme) else RGBColor(0xAA, 0xAA, 0xAA)
    box = slide.shapes.add_textbox(Inches(12.3), Inches(6.9), Inches(0.8), Inches(0.4))
    _text(box.text_frame, str(num), t["font_body"], Pt(10), pg_color, align=PP_ALIGN.RIGHT)


def _add_notes(slide, notes):
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _style_chart(chart, t, xl_type, theme_name):
    """Apply theme-aware styling to chart: axis colors, data labels, legend, per-point coloring."""
    is_dark = _is_dark(theme_name)
    axis_color = t["text_dark"]
    grid_color = RGBColor(0x44, 0x44, 0x55) if is_dark else RGBColor(0xE0, 0xE0, 0xE0)
    chart_colors = t.get("chart_colors", [t["accent"], t["primary"], t["accent2"]])

    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(11)
    chart.legend.font.name = t["font_body"]
    chart.legend.font.color.rgb = axis_color

    if xl_type in (XL_CHART_TYPE.COLUMN_CLUSTERED, XL_CHART_TYPE.LINE_MARKERS):
        cat_axis = chart.category_axis
        cat_axis.tick_labels.font.size = Pt(10)
        cat_axis.tick_labels.font.name = t["font_body"]
        cat_axis.tick_labels.font.color.rgb = axis_color
        cat_axis.format.line.color.rgb = grid_color

        val_axis = chart.value_axis
        val_axis.tick_labels.font.size = Pt(10)
        val_axis.tick_labels.font.name = t["font_body"]
        val_axis.tick_labels.font.color.rgb = axis_color
        val_axis.major_gridlines.format.line.color.rgb = grid_color
        val_axis.format.line.color.rgb = grid_color

    plot = chart.plots[0]

    if xl_type in (XL_CHART_TYPE.PIE, XL_CHART_TYPE.DOUGHNUT):
        for s_idx, series in enumerate(plot.series):
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = chart_colors[s_idx % len(chart_colors)]
            num_points = len(series.values)
            for pt_idx in range(num_points):
                point = series.points[pt_idx]
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = chart_colors[pt_idx % len(chart_colors)]

        plot.has_data_labels = True
        dl = plot.data_labels
        dl.font.size = Pt(11)
        dl.font.name = t["font_body"]
        dl.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        dl.font.bold = True
        dl.show_percentage = True
        dl.show_category_name = True
        dl.show_value = False
        dl.separator = "\n"

    elif xl_type == XL_CHART_TYPE.COLUMN_CLUSTERED:
        num_series = len(plot.series)
        if num_series == 1:
            series = plot.series[0]
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = chart_colors[0]
            num_points = len(series.values)
            for pt_idx in range(num_points):
                point = series.points[pt_idx]
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = chart_colors[pt_idx % len(chart_colors)]
            chart.has_legend = False
        else:
            for idx, series in enumerate(plot.series):
                fill = series.format.fill
                fill.solid()
                fill.fore_color.rgb = chart_colors[idx % len(chart_colors)]

        plot.has_data_labels = True
        dl = plot.data_labels
        dl.font.size = Pt(9)
        dl.font.name = t["font_body"]
        dl.font.color.rgb = axis_color
        dl.font.bold = True
        dl.show_value = True
        dl.show_category_name = False
        dl.show_series_name = False
        dl.number_format = '#,##0'
        dl.number_format_is_linked = False

    elif xl_type == XL_CHART_TYPE.LINE_MARKERS:
        for idx, series in enumerate(plot.series):
            fill = series.format.line
            fill.color.rgb = chart_colors[idx % len(chart_colors)]
            fill.width = Pt(2.5)
            marker = series.marker
            marker.style = 8  # circle
            marker.size = 8
            marker.format.fill.solid()
            marker.format.fill.fore_color.rgb = chart_colors[idx % len(chart_colors)]

        plot.has_data_labels = True
        dl = plot.data_labels
        dl.font.size = Pt(8)
        dl.font.name = t["font_body"]
        dl.font.color.rgb = axis_color
        dl.show_value = True
        dl.show_category_name = False
        dl.number_format = '#,##0'


# ---------------------------------------------------------------------------
# Layout builders
# ---------------------------------------------------------------------------

def _build_title_slide(prs, sd, theme):
    t = _t(theme)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, t["bg_dark"])

    img_path = fetch_image(sd.image_keyword) if sd.image_keyword else None
    if img_path:
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
        _rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, RGBColor(0, 0, 0), alpha=72)

    _rect(slide, Inches(0), Inches(1.2), SLIDE_WIDTH, Inches(3.6), t["bg_dark"], alpha=60)
    _rect(slide, Inches(0), Inches(3.2), SLIDE_WIDTH, Inches(0.08), t["accent"])

    box = slide.shapes.add_textbox(Inches(1.2), Inches(1.6), Inches(10.8), Inches(1.8))
    _text(box.text_frame, sd.title, t["font_title"], Pt(46),
          t["text_light"], bold=True, align=PP_ALIGN.CENTER)

    if sd.subtitle:
        _rect(slide, Inches(3), Inches(3.6), Inches(7.3), Inches(0.9), t["bg_dark"], alpha=50)
        box2 = slide.shapes.add_textbox(Inches(3.2), Inches(3.65), Inches(6.9), Inches(0.8))
        _text(box2.text_frame, sd.subtitle, t["font_body"], Pt(22),
              t["text_light"], align=PP_ALIGN.CENTER)

    _rect(slide, Inches(5.2), Inches(3.28), Inches(2.9), Inches(0.04), t["accent"])
    _add_notes(slide, sd.speaker_notes)


def _build_closing_slide(prs, sd, theme):
    t = _t(theme)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, t["bg_dark"])

    _rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_HEIGHT, t["accent"])
    _rect(slide, Inches(0), Inches(3.4), SLIDE_WIDTH, Inches(0.06), t["accent"])

    box = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.5))
    _text(box.text_frame, sd.title, t["font_title"], Pt(42),
          t["text_light"], bold=True, align=PP_ALIGN.CENTER)

    if sd.subtitle:
        box2 = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(9.3), Inches(0.8))
        _text(box2.text_frame, sd.subtitle, t["font_body"], Pt(20),
              t["text_light"], align=PP_ALIGN.CENTER)

    _add_notes(slide, sd.speaker_notes)


def _build_section_slide(prs, sd, theme):
    t = _t(theme)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, t["accent2"])

    _rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT, t["text_light"])

    box = slide.shapes.add_textbox(Inches(1.5), Inches(2.3), Inches(10.3), Inches(1.5))
    _text(box.text_frame, sd.title, t["font_title"], Pt(42),
          t["text_light"], bold=True)

    _rect(slide, Inches(1.5), Inches(4.1), Inches(3), Inches(0.06), t["text_light"])

    if sd.subtitle:
        box2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(10), Inches(0.8))
        _text(box2.text_frame, sd.subtitle, t["font_body"], Pt(20), t["text_light"])

    _add_notes(slide, sd.speaker_notes)


def _build_content_slide(prs, sd, theme):
    t = _t(theme)
    dark = _is_dark(theme)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, t["bg_light"])

    _rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), t["accent"])
    _rect(slide, Inches(0), Inches(0.06), Inches(0.08), Inches(1.2), t["primary"])

    tbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(11), Inches(0.9))
    _text(tbox.text_frame, sd.title, t["font_title"], Pt(28),
          t["text_light"] if dark else t["primary"], bold=True)

    _rect(slide, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.04), t["accent"])

    if sd.bullet_points:
        bx = slide.shapes.add_textbox(Inches(1.2), Inches(1.7), Inches(10.5), Inches(5.0))
        tf = bx.text_frame
        tf.word_wrap = True
        for i, pt in enumerate(sd.bullet_points):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"  {pt}"
            p.font.name = t["font_body"]
            p.font.size = Pt(17)
            p.font.color.rgb = t["text_dark"]
            p.space_before = Pt(10)
            p.space_after = Pt(10)

            dot = slide.shapes.add_textbox(
                Inches(0.75), Inches(1.75 + i * 0.55), Inches(0.3), Inches(0.35)
            )
            _text(dot.text_frame, "●", t["font_body"], Pt(9), t["accent"])

    _add_page_number(slide, sd.slide_number, theme)
    _add_notes(slide, sd.speaker_notes)


def _build_two_column_slide(prs, sd, theme):
    t = _t(theme)
    dark = _is_dark(theme)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, t["bg_light"])

    _rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), t["accent"])

    tbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(11), Inches(0.9))
    _text(tbox.text_frame, sd.title, t["font_title"], Pt(28),
          t["text_light"] if dark else t["primary"], bold=True)

    _rect(slide, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.04), t["accent"])

    mid = len(sd.bullet_points) // 2 if sd.bullet_points else 0
    left_pts = sd.bullet_points[:mid] if mid else []
    right_pts = sd.bullet_points[mid:] if mid else []

    for col_idx, points in enumerate([left_pts, right_pts]):
        x_off = Inches(0.8) if col_idx == 0 else Inches(7.0)
        if col_idx == 0:
            _rect(slide, Inches(6.55), Inches(1.6), Inches(0.03), Inches(5.0), t["accent"])

        bx = slide.shapes.add_textbox(x_off, Inches(1.7), Inches(5.5), Inches(5.0))
        tf = bx.text_frame
        tf.word_wrap = True
        for i, pt in enumerate(points):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"  {pt}"
            p.font.name = t["font_body"]
            p.font.size = Pt(16)
            p.font.color.rgb = t["text_dark"]
            p.space_before = Pt(8)
            p.space_after = Pt(8)

    _add_page_number(slide, sd.slide_number, theme)
    _add_notes(slide, sd.speaker_notes)


def _build_image_right_slide(prs, sd, theme):
    t = _t(theme)
    dark = _is_dark(theme)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, t["bg_light"])

    _rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), t["accent"])
    _rect(slide, Inches(0), Inches(0.06), Inches(0.08), Inches(1.2), t["primary"])

    tbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(6.5), Inches(0.9))
    _text(tbox.text_frame, sd.title, t["font_title"], Pt(28),
          t["text_light"] if dark else t["primary"], bold=True)

    _rect(slide, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.04), t["accent"])

    if sd.bullet_points:
        bx = slide.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(6.0), Inches(5.0))
        tf = bx.text_frame
        tf.word_wrap = True
        for i, pt in enumerate(sd.bullet_points):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"  {pt}"
            p.font.name = t["font_body"]
            p.font.size = Pt(16)
            p.font.color.rgb = t["text_dark"]
            p.space_before = Pt(8)
            p.space_after = Pt(8)

            dot = slide.shapes.add_textbox(
                Inches(0.65), Inches(1.75 + i * 0.5), Inches(0.3), Inches(0.35)
            )
            _text(dot.text_frame, "●", t["font_body"], Pt(8), t["accent"])

    img_path = fetch_image(sd.image_keyword) if sd.image_keyword else None
    img_left = Inches(7.4)
    img_top = Inches(1.0)
    img_w = Inches(5.5)
    img_h = Inches(5.8)

    if img_path:
        _rect(slide, img_left, img_top, img_w, img_h, t["primary"])
        slide.shapes.add_picture(
            img_path,
            img_left + Inches(0.1), img_top + Inches(0.1),
            img_w - Inches(0.2), img_h - Inches(0.2),
        )
    else:
        _rect(slide, img_left, img_top, img_w, img_h, t["bg_dark"])
        _rect(slide, img_left + Inches(1.5), img_top + Inches(1.8),
              Inches(2.5), Inches(2.5), t["accent"], alpha=40)
        _rect(slide, img_left + Inches(2.0), img_top + Inches(2.3),
              Inches(2.5), Inches(2.5), t["accent2"], alpha=30)

    _add_page_number(slide, sd.slide_number, theme)
    _add_notes(slide, sd.speaker_notes)


def _build_chart_slide(prs, sd, theme):
    t = _t(theme)
    dark = _is_dark(theme)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, t["bg_light"])

    _rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), t["accent"])
    _rect(slide, Inches(0), Inches(0.06), Inches(0.08), Inches(1.2), t["primary"])

    tbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(11), Inches(0.9))
    _text(tbox.text_frame, sd.title, t["font_title"], Pt(28),
          t["text_light"] if dark else t["primary"], bold=True)

    _rect(slide, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.04), t["accent"])

    if sd.chart and sd.chart.categories and sd.chart.series:
        chart_type_map = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
            "doughnut": XL_CHART_TYPE.DOUGHNUT,
        }
        xl_type = chart_type_map.get(sd.chart.chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

        chart_data = CategoryChartData()
        chart_data.categories = sd.chart.categories
        for s in sd.chart.series:
            chart_data.add_series(s.get("name", ""), s.get("values", []))

        chart_frame = slide.shapes.add_chart(
            xl_type,
            Inches(1.0), Inches(1.6), Inches(11.3), Inches(5.2),
            chart_data,
        )

        _style_chart(chart_frame.chart, t, xl_type, theme)
    else:
        box = slide.shapes.add_textbox(Inches(3), Inches(3), Inches(7), Inches(1))
        _text(box.text_frame, "（グラフデータなし）", t["font_body"], Pt(16),
              t["text_dark"], align=PP_ALIGN.CENTER)

    _add_page_number(slide, sd.slide_number, theme)
    _add_notes(slide, sd.speaker_notes)


def _build_stats_slide(prs, sd, theme):
    t = _t(theme)
    dark = _is_dark(theme)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, t["bg_light"])

    _rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), t["accent"])
    _rect(slide, Inches(0), Inches(0.06), Inches(0.08), Inches(1.2), t["primary"])

    tbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(11), Inches(0.9))
    _text(tbox.text_frame, sd.title, t["font_title"], Pt(28),
          t["text_light"] if dark else t["primary"], bold=True)

    _rect(slide, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.04), t["accent"])

    stats = sd.stats or []
    n = len(stats) if stats else 1
    card_w = min(Inches(3.0), (SLIDE_WIDTH - Inches(2.4)) / n - Inches(0.4))
    total_w = n * card_w + (n - 1) * Inches(0.5)
    start_x = (SLIDE_WIDTH - total_w) / 2

    for i, stat in enumerate(stats):
        x = start_x + i * (card_w + Inches(0.5))
        y = Inches(1.9)
        card_h = Inches(4.8)

        card_bg = RGBColor(0x22, 0x22, 0x40) if dark else RGBColor(0xFF, 0xFF, 0xFF)
        card = _rect(slide, x, y, card_w, card_h, card_bg)
        card.shadow.inherit = False
        _add_shadow(card, blur=Pt(10), offset=Pt(5))

        _rect(slide, x, y, card_w, Inches(0.10), t["accent"])

        color_idx = i % len(t.get("chart_colors", [t["accent"]]))
        indicator_color = t["chart_colors"][color_idx] if t.get("chart_colors") else t["accent"]
        _rect(slide, x + Inches(0.3), y + Inches(0.5),
              Inches(0.6), Inches(0.06), indicator_color)

        val_text = stat.value
        val_size = Pt(40)
        if len(val_text) > 10:
            val_size = Pt(32)
        elif len(val_text) > 7:
            val_size = Pt(36)

        val_box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(1.0),
                                           card_w - Inches(0.3), Inches(2.0))
        val_tf = val_box.text_frame
        val_tf.word_wrap = True
        val_p = val_tf.paragraphs[0]
        val_p.text = val_text
        val_p.font.name = t["font_title"]
        val_p.font.size = val_size
        val_p.font.color.rgb = indicator_color
        val_p.font.bold = True
        val_p.alignment = PP_ALIGN.CENTER
        val_box.text_frame.paragraphs[0].space_before = Pt(0)
        val_box.text_frame.paragraphs[0].space_after = Pt(0)

        lbl_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(3.2),
                                           card_w - Inches(0.2), Inches(1.3))
        lbl_tf = lbl_box.text_frame
        lbl_tf.word_wrap = True
        lbl_p = lbl_tf.paragraphs[0]
        lbl_p.text = stat.label
        lbl_p.font.name = t["font_body"]
        lbl_p.font.size = Pt(16)
        lbl_p.font.color.rgb = t["text_dark"]
        lbl_p.font.bold = False
        lbl_p.alignment = PP_ALIGN.CENTER

    _add_page_number(slide, sd.slide_number, theme)
    _add_notes(slide, sd.speaker_notes)


def _build_summary_slide(prs, summary_data: dict, theme: str):
    """Build a data analysis summary page from CSV/PDF analysis results."""
    t = _t(theme)
    dark = _is_dark(theme)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, t["bg_light"])

    _rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.3), t["bg_dark"])
    _rect(slide, Inches(0), Inches(1.3), SLIDE_WIDTH, Inches(0.06), t["accent"])

    title = summary_data.get("title", "データ分析サマリー")
    tbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(11), Inches(0.9))
    _text(tbox.text_frame, title, t["font_title"], Pt(30),
          t["text_light"], bold=True)

    source_type = summary_data.get("source_type", "")
    if source_type:
        sub = slide.shapes.add_textbox(Inches(9), Inches(0.3), Inches(4), Inches(0.5))
        label = "CSV分析" if source_type == "csv" else "PDF文書" if source_type == "pdf" else "CSV + PDF"
        _text(sub.text_frame, f"📋 {label}", t["font_body"], Pt(13),
              t["accent"], align=PP_ALIGN.RIGHT)

    sections = summary_data.get("sections", [])
    if not sections:
        return

    num_sections = len(sections)
    if num_sections <= 2:
        cols = 1
    elif num_sections <= 4:
        cols = 2
    else:
        cols = 3

    col_w = (SLIDE_WIDTH - Inches(1.6)) / cols - Inches(0.3)
    row_h = Inches(2.6) if num_sections <= 4 else Inches(2.4)
    start_x = Inches(0.8)
    start_y = Inches(1.65)

    for idx, section in enumerate(sections[:6]):
        col = idx % cols
        row = idx // cols
        x = start_x + col * (col_w + Inches(0.3))
        y = start_y + row * (row_h + Inches(0.2))

        sec_bg = RGBColor(0x22, 0x22, 0x40) if dark else RGBColor(0xFF, 0xFF, 0xFF)
        card = _rect(slide, x, y, col_w, row_h, sec_bg)
        _add_shadow(card, blur=Pt(6), offset=Pt(3))

        color_idx = idx % len(t.get("chart_colors", [t["accent"]]))
        bar_color = t["chart_colors"][color_idx] if t.get("chart_colors") else t["accent"]
        _rect(slide, x, y, Inches(0.08), row_h, bar_color)

        sec_title = section.get("heading", "")
        sec_tbox = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.15),
                                             col_w - Inches(0.4), Inches(0.5))
        _text(sec_tbox.text_frame, sec_title, t["font_title"], Pt(14),
              bar_color, bold=True)

        items = section.get("items", [])
        if items:
            items_box = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.65),
                                                  col_w - Inches(0.4), row_h - Inches(0.8))
            tf = items_box.text_frame
            tf.word_wrap = True
            for j, item in enumerate(items[:6]):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = f"• {item}"
                p.font.name = t["font_body"]
                p.font.size = Pt(11)
                p.font.color.rgb = t["text_dark"]
                p.space_before = Pt(2)
                p.space_after = Pt(2)

    _add_page_number(slide, 2, theme)


# ---------------------------------------------------------------------------
# Summary data builder
# ---------------------------------------------------------------------------

def build_summary_data(csv_analysis: dict | None = None,
                       pdf_text: str = "") -> dict | None:
    """Build structured summary data from CSV analysis and/or PDF text."""
    sections = []
    source_type = ""

    if csv_analysis and pdf_text:
        source_type = "both"
    elif csv_analysis:
        source_type = "csv"
    elif pdf_text:
        source_type = "pdf"
    else:
        return None

    if csv_analysis:
        overview = csv_analysis.get("overview", {})
        ov_items = [
            f"データ規模: {overview.get('rows', '?')}行 × {overview.get('columns', '?')}列",
        ]
        cols = overview.get("column_names", [])
        if cols:
            ov_items.append(f"列: {', '.join(cols[:8])}" + ("..." if len(cols) > 8 else ""))
        num_cols = overview.get("numeric_columns", [])
        cat_cols = overview.get("category_columns", [])
        if num_cols:
            ov_items.append(f"数値列: {', '.join(num_cols[:5])}")
        if cat_cols:
            ov_items.append(f"カテゴリ列: {', '.join(cat_cols[:5])}")
        sections.append({"heading": "データ概要", "items": ov_items})

        stats = csv_analysis.get("statistics", {})
        if stats:
            stat_items = []
            for col_name, s in list(stats.items())[:5]:
                mean = s.get("mean", "?")
                if isinstance(mean, float):
                    mean = f"{mean:,.2f}"
                stat_items.append(f"{col_name}: 平均 {mean} (最小 {s.get('min', '?')} 〜 最大 {s.get('max', '?')})")
            sections.append({"heading": "基本統計量", "items": stat_items})

        corrs = csv_analysis.get("correlations", [])
        if corrs:
            corr_items = []
            for c in corrs[:4]:
                val = c["value"]
                direction = "正の相関" if val > 0 else "負の相関"
                corr_items.append(f"{c['col1']} ↔ {c['col2']}: {val:.3f} ({direction})")
            sections.append({"heading": "注目すべき相関関係", "items": corr_items})

        trends = csv_analysis.get("trends", {})
        if trends:
            trend_items = [f"{col}: {desc}" for col, desc in list(trends.items())[:5]]
            sections.append({"heading": "トレンド分析", "items": trend_items})

        cat_summary = csv_analysis.get("category_summary", {})
        if cat_summary:
            for col_name, counts in list(cat_summary.items())[:2]:
                cat_items = [f"{k}: {v}件" for k, v in list(counts.items())[:6]]
                sections.append({"heading": f"{col_name}別 分布", "items": cat_items})

    if pdf_text:
        lines = pdf_text.strip().split("\n")
        meaningful = [l.strip() for l in lines if len(l.strip()) > 5]
        preview = meaningful[:8]
        pdf_items = [l[:80] + ("..." if len(l) > 80 else "") for l in preview]
        sections.append({"heading": "PDF文書の概要", "items": pdf_items})

    if not sections:
        return None

    return {
        "title": "データ分析サマリー",
        "source_type": source_type,
        "sections": sections,
    }


# ---------------------------------------------------------------------------
# Main builder
# ---------------------------------------------------------------------------

LAYOUT_BUILDERS = {
    "title_slide": _build_title_slide,
    "closing": _build_closing_slide,
    "section_header": _build_section_slide,
    "content": _build_content_slide,
    "two_column": _build_two_column_slide,
    "image_right": _build_image_right_slide,
    "chart": _build_chart_slide,
    "stats": _build_stats_slide,
}


def build_pptx(data: PresentationData, filename: str,
               summary_data: dict | None = None) -> str:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    theme = data.theme if data.theme in THEMES else "ocean_blue"

    for i, sd in enumerate(data.slides):
        builder = LAYOUT_BUILDERS.get(sd.layout, _build_content_slide)
        builder(prs, sd, theme)

        if i == 0 and summary_data:
            _build_summary_slide(prs, summary_data, theme)

    safe = "".join(
        c for c in filename
        if c.isalnum() or c in (" ", "_", "-", ".", "ー")
        or ("\u3000" <= c <= "\u9fff") or ("\uac00" <= c <= "\ud7af")
    ).strip() or "presentation"
    path = os.path.join(OUTPUT_DIR, f"{safe}.pptx")

    prs.save(path)
    return path
